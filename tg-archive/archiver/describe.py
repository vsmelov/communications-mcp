"""Описание картинок и PDF дешёвой vision-моделью (OpenRouter) + журнал трат.

Идея как у транскрибации аудио: заменяем бинарь текстом.
  - картинка -> модель описывает содержимое и дословно приводит видимый текст;
  - PDF      -> текст извлекается ЛОКАЛЬНО и бесплатно (pypdf); скан без текста
               описывается моделью как картинка первой страницы (если получится),
               иначе файл остаётся на диске.
После успеха файл удаляется, строка помечается `described`; оригинал докачиваем
по msg_id (лежит на серверах Telegram). Управляется секцией describe в config.

Траты идут в data/describe_ledger.jsonl (тот же формат, что у транскрибации):
в /status и на дашборде видно всего / за сутки / за час, есть дневной стоп-кран.
"""
from __future__ import annotations

import asyncio
import base64
import logging
import mimetypes
from pathlib import Path

from .scrub import SecretVault, scrub_text

log = logging.getLogger("describe")

# $ за 1M токенов (вход, выход). Сверять с openrouter.ai/models.
PRICING = {
    "google/gemini-2.0-flash-lite-001": (0.075, 0.30),
    "google/gemini-2.5-flash-lite": (0.10, 0.40),
    "google/gemini-flash-1.5-8b": (0.038, 0.15),
}
_DEFAULT_PRICE = (0.10, 0.40)

_IMAGE_EXTS = (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".heic", ".heif")
# Документы, из которых текст извлекается ЛОКАЛЬНО и бесплатно (без модели).
_DOC_EXTS = (".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt", ".md")

# English на выходе — заметно дешевле по токенам, чем кириллица. Дословный
# текст (OCR) при этом остаётся в оригинальном языке. Переопределяется в
# config.yaml -> describe.prompt.
DEFAULT_PROMPT = (
    "Analyze this image and answer in English, as labeled fields. Do not invent "
    "anything not visible; if a field does not apply, write \"none\".\n"
    "Type: photo / screenshot / chart / document scan / diagram / meme / selfie / other\n"
    "Description: what is shown — objects, people, scene, action (2-4 sentences)\n"
    "Text: ALL visible text, VERBATIM, in its ORIGINAL language (keep numbers, "
    "contract addresses, tickers, links, dates exactly). If none — \"none\".\n"
    "Text language: language(s) of the text above\n"
    "Data: key numbers / amounts / dates / contract addresses (0x…) / tickers / metrics\n"
    "People: how many and what they are doing (do not guess names)\n"
    "Mood: tone or mood if relevant (meme, celebration, tense, etc.)\n"
    "Quality: clear / blurry / phone-screen photo / low-res / cropped\n"
    "Sensitive: \"yes\" + what, if it shows secrets (passport, bank card, private key, "
    "seed phrase, password); otherwise \"no\"\n"
    "Context: if a screenshot of an app / chat / exchange / website — which and about what\n"
    "Notes: anything unusual"
)

# Документы описываются ПОДРОБНЕЕ картинок: это саммери сути, по которому
# агент решит, надо ли докачивать оригинал. Оригинал в Telegram (докачиваем по id).
DEFAULT_DOC_PROMPT = (
    "Below is text extracted from a document file. Write a DETAILED summary in "
    "English so a reader understands what the document is and contains without "
    "opening it. Be thorough. Do not invent; if a field is unknown, write \"unknown\".\n"
    "Type: contract / whitepaper / pitch deck / report / spreadsheet / invoice / "
    "legal / tokenomics / spec / resume / other\n"
    "Title/subject: the document's title or main subject\n"
    "Summary: a thorough summary of the content — main points, structure, sections, "
    "conclusions (6-12 sentences)\n"
    "Key data: important numbers, amounts, dates, names, wallet/contract addresses, "
    "tickers, percentages, terms — as a list\n"
    "Parties/authors: who made it / who it is for / mentioned entities\n"
    "Purpose: what this document is for and its likely context\n"
    "Language: language(s) of the document\n"
    "Sensitive: \"yes\" + what, if it contains secrets or personal IDs; otherwise \"no\""
)


class Describer:
    def __init__(self, cfg, ledger, api_key: str | None, vault_path=None):
        self.cfg = cfg  # DescribeCfg
        self.ledger = ledger
        self.api_key = api_key
        self._client = None
        self._vault_path = vault_path   # data/secrets.sqlite3 (общий с хостовым CLI)
        self._sv = None

    def _vault(self):
        """Ленивое хранилище секретов. Одно на процесс, пишется из рабочих
        потоков describe последовательно (SecretVault: check_same_thread=False)."""
        if self._sv is None and self._vault_path is not None:
            self._sv = SecretVault(self._vault_path)
        return self._sv

    def available(self) -> bool:
        return bool(self.cfg.enabled and self.api_key)

    def _openai(self):
        if self._client is None:
            from openai import OpenAI

            self._client = OpenAI(
                base_url="https://openrouter.ai/api/v1", api_key=self.api_key
            )
        return self._client

    # ------------------------------------------------------------- отбор --
    def is_target_kind(self, kind: str | None) -> bool:
        return kind in self.cfg.kinds

    def _describable_file(self, row: dict, store) -> Path | None:
        """Путь к файлу на диске, если его вообще стоит описывать (картинка/PDF)."""
        f = row.get("file")
        if not f:
            return None
        p = store.dir / f
        if not p.is_file() or p.stat().st_size == 0:  # 0 байт = битая загрузка
            return None
        if p.suffix.lower() in _IMAGE_EXTS or p.suffix.lower() in _DOC_EXTS:
            return p
        return None

    def needs(self, row: dict, store) -> bool:
        if not self.available():
            return False
        if not self.is_target_kind(row.get("media")):
            return False
        if row.get("description") or row.get("described"):
            return False
        if row.get("desc_attempts", 0) >= self.cfg.error_retries:
            return False
        size = row.get("size_bytes") or 0
        if size and size > self.cfg.max_bytes:
            return False
        return True

    # ------------------------------------------------------------- проход --
    async def pass_for(self, syncer) -> dict:
        """Описать до per_cycle картинок/PDF чата, не выходя за дневной бюджет.

        Работает по УЖЕ скачанным файлам (media-фаза их кладёт на диск).
        Картинки/PDF -> текст -> файл удаляется. Прочие документы (zip и т.п.)
        не трогаются.
        """
        if not self.available():
            return {}
        store = syncer.store
        day_spent = self.ledger.spent()["day_usd"]
        pending = []
        for r in store.sorted_rows():
            if not self.needs(r, store):
                continue
            if self._describable_file(r, store) is not None:
                pending.append(r)
            if len(pending) >= self.cfg.per_cycle:
                break
        if not pending:
            return {}

        done = failed = 0
        usd_sum = 0.0
        for row in pending:
            if day_spent >= self.cfg.daily_budget_usd:
                log.warning("[%s] дневной бюджет описания $%.2f исчерпан",
                            store.slug, self.cfg.daily_budget_usd)
                break
            path = self._describable_file(row, store)
            if path is None:
                continue
            await syncer._sleep(self.cfg.delay_sec)
            is_doc = path.suffix.lower() in _DOC_EXTS  # текст извлекаем локально, суть — моделью
            ref = f"{store.chat_id}:{row['id']}"
            try:
                if is_doc:
                    text, usd, ptok, ctok, secret_kinds = await asyncio.to_thread(
                        self._describe_document, path, ref)
                else:
                    text, usd, ptok, ctok, secret_kinds = await asyncio.to_thread(
                        self._describe_image, path, ref)
                if not text:
                    row["desc_attempts"] = row.get("desc_attempts", 0) + 1
                    failed += 1
                    continue
                row["description"] = text[: self.cfg.max_text_chars]
                if secret_kinds:
                    # в документе (вход) или в описании картинки (выход) нашли секрет
                    row["media_secret"] = secret_kinds
                row["description_kind"] = "document" if is_doc else "image"
                row["description_model"] = self.cfg.model
                row.pop("description_error", None)
                self.ledger.add(
                    task="describe", kind=(path.suffix.lower().lstrip(".") if is_doc else "image"),
                    provider=self.cfg.provider, chat_id=store.chat_id, msg_id=row["id"],
                    model=self.cfg.model, prompt_tokens=ptok, completion_tokens=ctok, usd=usd,
                )
                day_spent += usd
                usd_sum += usd
                # заменили бинарь текстом — файл больше не нужен
                if self.cfg.drop_after:
                    try:
                        path.unlink(missing_ok=True)
                    except Exception:
                        pass
                    row.pop("file", None)
                    row["described"] = True   # докачиваемо по msg_id
                done += 1
            except Exception as exc:
                row["desc_attempts"] = row.get("desc_attempts", 0) + 1
                row["description_error"] = str(exc)[:200]
                failed += 1
                log.warning("[%s] describe %s: %s", store.slug, row["id"], exc)
        return {"described": done, "failed": failed, "usd": round(usd_sum, 4)}

    # ------------------------------------------------------------- модели --
    def _price(self):
        return PRICING.get(self.cfg.model, _DEFAULT_PRICE)

    @staticmethod
    def _sniff_mime(data: bytes) -> str | None:
        """Тип картинки по магическим байтам (расширение у телеграм-файлов врёт)."""
        if data[:3] == b"\xff\xd8\xff":
            return "image/jpeg"
        if data[:8] == b"\x89PNG\r\n\x1a\n":
            return "image/png"
        if data[:6] in (b"GIF87a", b"GIF89a"):
            return "image/gif"
        if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
            return "image/webp"
        if data[4:12] in (b"ftypheic", b"ftypheix", b"ftypmif1", b"ftyphevc"):
            return "image/heic"
        return None

    def _describe_image(self, path: Path, ref: str | None = None
                        ) -> tuple[str, float, int, int, list]:
        raw = path.read_bytes()
        mime = self._sniff_mime(raw)
        if mime is None:
            # не распознанная как картинка — не шлём в модель (битый/чужой формат)
            raise ValueError(f"не картинка по сигнатуре: {path.name}")
        b64 = base64.b64encode(raw).decode()
        prompt = (self.cfg.prompt or "").strip() or DEFAULT_PROMPT
        resp = self._openai().chat.completions.create(
            model=self.cfg.model,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url",
                     "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }],
            max_tokens=1000,
        )
        text = (resp.choices[0].message.content or "").strip()
        usd, ptok, ctok = self._cost(resp)
        # Картинку в модель уже отправили (иначе не описать), но если это был
        # СКРИНШОТ ключа/сида и модель переписала секрет в описание — вычищаем
        # его из сохраняемого текста и помечаем строку. Оригинал не хранится.
        text, counts = scrub_text(text, self._vault(), ref)
        return text, usd, ptok, ctok, sorted(counts)

    def _describe_document(self, path: Path, ref: str | None = None
                           ) -> tuple[str, float, int, int, list]:
        """Документ -> подробное саммери сути. Текст извлекаем локально
        (бесплатно), затем прогоняем через модель для детального структурного
        резюме (подробнее, чем описание картинки). Оригинал не хранится, но
        докачиваем по msg_id — в саммери это указано рендером. Если текста нет
        (скан PDF/пустой) — возвращаем "" (пропуск)."""
        raw = self._extract_text(path)
        if len(raw) < 20:
            return "", 0.0, 0, 0, []
        prompt = (self.cfg.doc_prompt or "").strip() or DEFAULT_DOC_PROMPT
        clipped = raw[: self.cfg.doc_max_input_chars]
        # ВЫРЕЗАЕМ секреты из извлечённого текста ДО отправки в модель
        clipped, counts = scrub_text(clipped, self._vault(), ref)
        resp = self._openai().chat.completions.create(
            model=self.cfg.model,
            messages=[{
                "role": "user",
                "content": f"{prompt}\n\nFile name: {path.name}\n\n--- extracted text ---\n{clipped}",
            }],
            max_tokens=1200,
        )
        summary = (resp.choices[0].message.content or "").strip()
        usd, ptok, ctok = self._cost(resp)
        return summary, usd, ptok, ctok, sorted(counts)

    def _extract_text(self, path: Path) -> str:
        """Локальная бесплатная выемка текста из pdf/docx/xlsx/pptx/csv/txt."""
        ext = path.suffix.lower()
        text = ""
        try:
            if ext == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                text = "\n".join((pg.extract_text() or "") for pg in reader.pages[:50]).strip()
            elif ext == ".docx":
                from docx import Document
                doc = Document(str(path))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for tbl in doc.tables:
                    for r in tbl.rows:
                        parts.append(" | ".join(c.text for c in r.cells))
                text = "\n".join(parts).strip()
            elif ext == ".xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(str(path), read_only=True, data_only=True)
                lines = []
                for ws in wb.worksheets:
                    lines.append(f"[sheet: {ws.title}]")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None]
                        if cells:
                            lines.append(" | ".join(cells))
                        if len(lines) > 2000:
                            break
                wb.close()
                text = "\n".join(lines).strip()
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[slide {i}]")
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            parts.append(shape.text_frame.text)
                text = "\n".join(parts).strip()
            elif ext in (".csv", ".txt", ".md"):
                text = path.read_text(encoding="utf-8", errors="replace").strip()
        except Exception as exc:
            log.warning("выемка текста не удалась %s: %s", path.name, exc)
        return text

    def _cost(self, resp) -> tuple[float, int, int]:
        u = getattr(resp, "usage", None)
        if not u:
            return 0.0, 0, 0
        ptok = getattr(u, "prompt_tokens", 0) or 0
        ctok = getattr(u, "completion_tokens", 0) or 0
        pin, pout = self._price()
        return ptok / 1e6 * pin + ctok / 1e6 * pout, ptok, ctok
